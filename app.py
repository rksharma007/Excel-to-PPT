from flask import Flask, request, render_template, send_file
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import random
import os
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LABEL_POSITION

app = Flask(__name__)

# Custom color palette
# custom_colors = [
#     '#4B8BBE',  # Blue
#     '#FFCC00',  # Yellow
#     '#5DADEC',  # Light Blue
#     '#9B59B6',  # Purple
#     '#E74C3C',  # Red
#     '#2ECC71',  # Green
#     '#F39C12',  # Orange
#     '#34495E',  # Dark Blue
#     '#95A5A6',  # Gray
#     '#1ABC9C',  # Teal
#     '#C0392B',  # Dark Red
#     '#8E44AD',  # Dark Purple
#     '#16A085',  # Dark Teal
#     '#27AE60',  # Forest Green
#     '#2980B9',  # Ocean Blue
#     '#D35400',  # Burnt Orange
#     '#7F8C8D',  # Cool Gray
#     '#BDC3C7',  # Light Gray
#     '#F1C40F',  # Golden Yellow
#     '#E67E22',  # Pumpkin Orange
#     '#3498DB',  # Sky Blue
#     '#9C640C',  # Brownish Yellow
#     '#7D3C98',  # Deep Purple
#     '#F1948A',  # Light Coral
#     '#76D7C4',  # Aqua Green
#     '#F8C471',  # Sand Yellow
#     '#ABB2B9',  # Steel Gray
#     '#48C9B0',  # Mint Green
#     '#F5B041',  # Bright Orange
#     '#D98880',  # Rosy Red
# ]

custom_colors = [
    '#4B8BBE', '#FFCC00', '#5DADEC', '#9B59B6', '#E74C3C',
    '#2ECC71', '#F39C12', '#34495E', '#95A5A6', '#1ABC9C',
    '#C0392B', '#8E44AD', '#16A085', '#27AE60', '#2980B9',
    '#D35400', '#7F8C8D', '#BDC3C7', '#F1C40F', '#E67E22',
    '#3498DB', '#9C640C', '#7D3C98', '#F1948A', '#76D7C4',
    '#F8C471', '#ABB2B9', '#48C9B0', '#F5B041', '#D98880'
]



def create_presentation(file_path, selected_colors, create_table, create_pie_chart):
    data = pd.read_excel(file_path, sheet_name=0)
    prs = Presentation()

    slide_number = 1

    for column in data.select_dtypes(include='object').columns:
        counts = data[column].value_counts(dropna=True)

        if counts.empty:
            continue

        counts.index = counts.index.astype(str)
        total_responses = len(data)

        if not create_pie_chart and not create_table:
            table_title = f"Table {slide_number}. Distribution on the basis of {column} (n={total_responses})"
            create_table_slide(prs, table_title, counts.index.tolist(), counts.values.astype(int))
            figure_title = f"Figure {slide_number}.  Distribution on the basis of {column} (n={total_responses})"
            create_chart_slide(prs, figure_title, counts.index.tolist(), counts.values.astype(int), selected_colors)

        if create_table:
            table_title = f"Table {slide_number}. Distribution on the basis of {column} (n={total_responses})"
            create_table_slide(prs, table_title, counts.index.tolist(), counts.values.astype(int))

        if create_pie_chart:
            figure_title = f"Figure {slide_number}.  Distribution on the basis of {column} (n={total_responses})"
            create_chart_slide(prs, figure_title, counts.index.tolist(), counts.values.astype(int), selected_colors)

        slide_number += 1

    output_file = 'Generated-PPT-RK.pptx'
    prs.save(output_file)
    return output_file

def create_table_slide(prs, title, labels, sizes):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

    # Add title
    title_box = slide.shapes.title
    title_box.text = title

    # Set title font size
    for paragraph in title_box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(20)  # Set font size to 20
            run.font.color.rgb = RGBColor(0, 0, 0)  # Optional: set font color to black

    # Create a table for frequency and percentage
    rows = len(labels) + 2  # +2 for header row and total row
    cols = 3  # Frequency, Percentage

    left = Inches(1)  # Centering will be done later
    top = Inches(1.5)
    width = Inches(5)
    height = Inches(2 + rows * 0.5)  # Adjust height based on number of rows

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Adjust the height of each row
    for row in table.rows:
        row.height = Inches(0.5)  # Set the desired height for each row

    # Set column names
    table.cell(0, 0).text = 'Category'
    table.cell(0, 1).text = 'Frequency'
    table.cell(0, 2).text = 'Percentage'

    total_responses = len(labels)
    total_frequency = sum(sizes)

    for i in range(len(labels)):
        table.cell(i + 1, 0).text = labels[i]
        freq = int(sizes[i]) if sizes[i] is not None else 0  # Ensure frequency is an integer
        perc = (freq / total_responses) * 100 if total_responses > 0 else 0
        table.cell(i + 1, 1).text = str(freq)
        table.cell(i + 1, 2).text = f"{perc:.1f}%"

    # Add Total row
    total_row_index = len(labels) + 1
    table.cell(total_row_index, 0).text = 'Total'
    table.cell(total_row_index, 1).text = str(total_frequency)

    total_percentage = (total_frequency / total_responses) * 100 if total_responses > 0 else 0
    table.cell(total_row_index, 2).text = f"{total_percentage:.1f}%"

    # Set background color for Total row
    for j in range(cols):
        cell = table.cell(total_row_index, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(79, 129, 189)   # Set background color to #4F81BD

        # Set font size and bold for Total row
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)   # Set font size to 18
                run.font.bold = True      # Make font bold
                run.font.color.rgb = RGBColor(255, 255, 255)

    # Center the table horizontally on the slide
    slide_width = prs.slide_width
    table_width = width
    left_centered = (slide_width - table_width) / 2
    for shape in slide.shapes:
        if shape.has_table:
            shape.left = int(left_centered)

def create_chart_slide(prs, title, labels, sizes, selected_colors):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

    # Add title
    title_box = slide.shapes.title
    title_box.text = title

    # Set title font size
    for paragraph in title_box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(20)  # Set font size to 20
            run.font.color.rgb = RGBColor(0, 0, 0)  # Optional: set font color to black

    # Create pie chart data
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series('', sizes)  # Add a series for the pie chart

    # Add the pie chart to the slide
    x, y, cx, cy = Inches(2.5), Inches(2), Inches(5), Inches(4)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

    # Set colors for the pie chart slices
    if selected_colors and len(selected_colors) >= len(labels):
        for i in range(len(labels)):
            chart.series[0].points[i].format.fill.solid()
            chart.series[0].points[i].format.fill.fore_color.rgb = RGBColor(
                int(selected_colors[i][1:3], 16),
                int(selected_colors[i][3:5], 16),
                int(selected_colors[i][5:7], 16)
            )
    else:
        # Fallback to default colors if not enough selected colors
        default_colors = ['#4B8BBE', '#FFCC00', '#5DADEC', '#9B59B6']
        for i in range(len(labels)):
            chart.series[0].points[i].format.fill.solid()
            color_index = i % len(default_colors)
            chart.series[0].points[i].format.fill.fore_color.rgb = RGBColor(
                int(default_colors[color_index][1:3], 16),
                int(default_colors[color_index][3:5], 16),
                int(default_colors[color_index][5:7], 16)
            )

    # Set legend options
    chart.has_legend = True
    chart.legend.include_in_layout = False

    # Add data labels with percentages and category names
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.show_value = False
    data_labels.show_percentage = True

    data_labels.number_format = '0.00%'
    data_labels.position = XL_LABEL_POSITION.BEST_FIT
    
# Define the main route to handle file uploads and presentation generation.
@app.route('/', methods=['GET', 'POST'])
def index():
   if request.method == 'POST':
       file = request.files['file']
       colors_selected = request.form.getlist('colors')
       include_table = 'table' in request.form.getlist('charts')
       include_pie_chart = 'pie' in request.form.getlist('charts')

       # Save the uploaded file temporarily.
       file_path_temp = os.path.join('uploads', file.filename)
       file.save(file_path_temp)

       # Create presentation based on user selections.
       output_file_path = create_presentation(file_path_temp,
                                               colors_selected,
                                               include_table,
                                               include_pie_chart)

       return send_file(output_file_path,
                        as_attachment=True,
                        download_name='ExceltoPPT-RK.pptx')

   return render_template('index.html')

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))  # Default to 5000 if PORT is not set
    app.run(host="0.0.0.0", port=port, debug=True)